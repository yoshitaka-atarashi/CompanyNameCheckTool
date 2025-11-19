from flask import Flask, render_template, request, jsonify, send_file
import os
from werkzeug.utils import secure_filename
from pptx import Presentation
from pptx.util import Pt
import re
import json
from io import BytesIO

app = Flask(__name__)

# 設定ファイルを読み込む
def load_config():
    """設定ファイルを読み込む"""
    config_file = 'config.json'
    default_config = {
        'default_keywords': ['Hitachi Astemo', '日立Astemo', '日立アステモ'],
        'default_replacement': 'Astemo',
        'max_file_size_mb': 50,
        'allowed_extensions': ['pptx', 'ppt']
    }
    
    if os.path.exists(config_file):
        try:
            with open(config_file, 'r', encoding='utf-8') as f:
                return json.load(f)
        except Exception as e:
            print(f"設定ファイルの読み込みに失敗しました: {str(e)}")
            return default_config
    return default_config

config = load_config()

# 設定
UPLOAD_FOLDER = 'uploads'
ALLOWED_EXTENSIONS = set(config['allowed_extensions'])
MAX_FILE_SIZE = config['max_file_size_mb'] * 1024 * 1024

if not os.path.exists(UPLOAD_FOLDER):
    os.makedirs(UPLOAD_FOLDER)

app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER
app.config['MAX_CONTENT_LENGTH'] = MAX_FILE_SIZE


def allowed_file(filename):
    """ファイルが許可されている拡張子かチェック"""
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS


def cleanup_uploads(files_to_delete):
    """処理後にアップロードフォルダをクリーンアップ"""
    try:
        for file_path in files_to_delete:
            if os.path.isfile(file_path):
                os.remove(file_path)
                print(f"クリーンアップ: {file_path}")
            elif os.path.isdir(file_path):
                # ディレクトリの場合は再帰的に削除
                import shutil
                shutil.rmtree(file_path, ignore_errors=True)
                print(f"クリーンアップ（フォルダ）: {file_path}")
    except Exception as e:
        print(f"クリーンアップエラー: {str(e)}")


def get_all_files_recursive(file_path, recursive=False):
    """ファイルリストを取得（再帰的に処理する場合）"""
    files_list = []
    
    try:
        if os.path.isfile(file_path):
            files_list.append(file_path)
        elif os.path.isdir(file_path) and recursive:
            # ディレクトリ内のすべてのファイルを再帰的に取得
            for root, dirs, files in os.walk(file_path):
                for file in files:
                    if allowed_file(file):
                        files_list.append(os.path.join(root, file))
    except Exception as e:
        print(f"ファイルリスト取得エラー: {str(e)}")
    
    return files_list


def find_keywords_in_presentation(prs, keywords):
    """プレゼンテーション内のキーワードを検出 (OR条件)
    通常スライドとマスタースライドの両方をチェック"""
    results = []
    
    # 通常スライドを処理
    for slide_num, slide in enumerate(prs.slides, 1):
        for shape_num, shape in enumerate(slide.shapes):
            if hasattr(shape, "text") and shape.text.strip():
                found_keywords = []
                total_count = 0
                
                # すべてのキーワードを検查
                for keyword in keywords:
                    if keyword.lower() in shape.text.lower():
                        count = shape.text.lower().count(keyword.lower())
                        found_keywords.append(keyword)
                        total_count += count
                
                # いずれかのキーワードが見つかった場合
                if found_keywords:
                    results.append({
                        'slide': slide_num,
                        'shape': shape_num,
                        'text': shape.text,
                        'keywords': found_keywords,
                        'count': total_count,
                        'is_master': False
                    })
    
    # マスタースライドを処理
    try:
        for master_num, master in enumerate(prs.slide_master.slide_layouts):
            for shape_num, shape in enumerate(master.shapes):
                if hasattr(shape, "text") and shape.text.strip():
                    found_keywords = []
                    total_count = 0
                    
                    for keyword in keywords:
                        if keyword.lower() in shape.text.lower():
                            count = shape.text.lower().count(keyword.lower())
                            found_keywords.append(keyword)
                            total_count += count
                    
                    if found_keywords:
                        results.append({
                            'slide': f'Master {master_num + 1}',
                            'shape': shape_num,
                            'text': shape.text,
                            'keywords': found_keywords,
                            'count': total_count,
                            'is_master': True
                        })
    except Exception as e:
        print(f"マスタースライド処理エラー: {str(e)}")
    
    return results


def replace_text_in_shape(shape, keywords, new_text, is_delete=False):
    """シェイプ内のテキストを置換 (複数キーワード対応)"""
    if hasattr(shape, "text_frame"):
        for paragraph in shape.text_frame.paragraphs:
            # パラグラフレベルでテキスト全体を取得
            full_text = ''.join(run.text for run in paragraph.runs)
            
            # いずれかのキーワードがマッチするか確認
            has_match = any(keyword.lower() in full_text.lower() for keyword in keywords)
            
            if not has_match:
                continue
            
            # すべてのキーワードを置換（ループで処理）
            new_full_text = full_text
            for keyword in keywords:
                pattern = re.compile(re.escape(keyword), re.IGNORECASE)
                new_full_text = pattern.sub(new_text, new_full_text)
            
            # すべてのrunをクリアして新しいテキストを設定
            for run in paragraph.runs:
                run.text = ''
            
            # 新しいテキストを最初のrunに設定
            if paragraph.runs:
                paragraph.runs[0].text = new_full_text
            else:
                # runがない場合は新しく作成
                paragraph.text = new_full_text


def process_presentation(prs, keywords, new_keyword=None, is_delete=False):
    """プレゼンテーション全体を処理 (複数キーワード対応)
    通常スライドとマスタースライドの両方を処理"""
    modified_count = 0
    
    # 通常スライドを処理
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text_frame"):
                original_text = shape.text
                
                # いずれかのキーワードが含まれているか確認
                has_keyword = any(kw.lower() in original_text.lower() for kw in keywords)
                
                if has_keyword:
                    # 置換先のテキストを決定
                    replacement_text = '' if is_delete else (new_keyword or keywords[0])
                    
                    # 置換を実行
                    replace_text_in_shape(shape, keywords, replacement_text, is_delete=is_delete)
                    
                    # 変更があったかを確認
                    if shape.text != original_text:
                        modified_count += 1
    
    # マスタースライドを処理
    try:
        for master in prs.slide_master.slide_layouts:
            for shape in master.shapes:
                if hasattr(shape, "text_frame"):
                    original_text = shape.text
                    
                    has_keyword = any(kw.lower() in original_text.lower() for kw in keywords)
                    
                    if has_keyword:
                        replacement_text = '' if is_delete else (new_keyword or keywords[0])
                        replace_text_in_shape(shape, keywords, replacement_text, is_delete=is_delete)
                        
                        if shape.text != original_text:
                            modified_count += 1
    except Exception as e:
        print(f"マスタースライド処理エラー: {str(e)}")
    
    return modified_count


@app.route('/')
def index():
    """メインページ"""
    return render_template('index.html', 
                          default_keywords=config['default_keywords'],
                          default_replacement=config['default_replacement'])


@app.route('/api/detect', methods=['POST'])
def detect_keywords():
    """キーワード検出API"""
    files_to_cleanup = []
    try:
        files = request.files.getlist('file')  # 複数ファイルに対応
        if not files or (len(files) == 1 and files[0].filename == ''):
            return jsonify({'error': 'ファイルがアップロードされていません'}), 400
        
        keywords_json = request.form.get('keywords', '[]')
        recursive = request.form.get('recursive', 'false').lower() == 'true'
        
        try:
            keywords = json.loads(keywords_json)
        except json.JSONDecodeError:
            keywords = [keywords_json] if keywords_json else []
        
        if not keywords or len(keywords) == 0:
            return jsonify({'error': 'キーワードを入力してください'}), 400
        
        # 複数ファイルを処理
        files_to_process = []
        for file in files:
            if file.filename == '':
                continue
            # 隠しファイルを除外
            if file.filename.startswith('.'):
                print(f"スキップ（隠しファイル）: {file.filename}")
                continue
            if not allowed_file(file.filename):
                print(f"スキップ（拡張子不可）: {file.filename}")
                continue
            
            filename = secure_filename(file.filename)
            filepath = os.path.join(app.config['UPLOAD_FOLDER'], filename)
            file.save(filepath)
            files_to_cleanup.append(filepath)
            files_to_process.append(filepath)
            print(f"処理対象に追加: {filename}")
        
        if not files_to_process:
            cleanup_uploads(files_to_cleanup)
            return jsonify({'error': '処理するPPTXファイルが見つかりません'}), 400
        
        # 全ファイルの結果を集約
        all_results = []
        total_count = 0
        total_affected_slides = 0
        
        for file_path in files_to_process:
            try:
                prs = Presentation(file_path)
                results = find_keywords_in_presentation(prs, keywords)
                
                # ファイル情報を結果に追加
                for result in results:
                    result['file'] = os.path.basename(file_path)
                
                all_results.extend(results)
                total_count += sum(r['count'] for r in results)
                total_affected_slides += len(results)
            except Exception as e:
                print(f"ファイル処理エラー {file_path}: {str(e)}")
                continue
        
        response = jsonify({
            'success': True,
            'keywords': keywords,
            'total_count': total_count,
            'affected_slides': total_affected_slides,
            'files_processed': len(files_to_process),
            'results': all_results
        })
        
        # 処理後にクリーンアップ
        cleanup_uploads(files_to_cleanup)
        return response
    
    except Exception as e:
        cleanup_uploads(files_to_cleanup)
        return jsonify({'error': f'エラーが発生しました: {str(e)}'}), 500


@app.route('/api/replace', methods=['POST'])
def replace_keywords():
    """キーワード置換API"""
    files_to_cleanup = []
    try:
        files = request.files.getlist('file')  # 複数ファイルに対応
        if not files or (len(files) == 1 and files[0].filename == ''):
            return jsonify({'error': 'ファイルがアップロードされていません'}), 400
        
        keywords_json = request.form.get('keywords', '[]')
        new_keyword = request.form.get('new_keyword', '').strip()
        action = request.form.get('action', 'replace')
        recursive = request.form.get('recursive', 'false').lower() == 'true'
        
        try:
            keywords = json.loads(keywords_json)
        except json.JSONDecodeError:
            keywords = [keywords_json] if keywords_json else []
        
        if not keywords or len(keywords) == 0:
            return jsonify({'error': 'キーワードを入力してください'}), 400
        
        if action == 'replace' and not new_keyword:
            return jsonify({'error': '置換先のキーワードを入力してください'}), 400
        
        # 複数ファイルを処理
        files_to_process = []
        for file in files:
            if file.filename == '':
                continue
            # 隠しファイルを除外
            if file.filename.startswith('.'):
                print(f"スキップ（隠しファイル）: {file.filename}")
                continue
            if not allowed_file(file.filename):
                print(f"スキップ（拡張子不可）: {file.filename}")
                continue
            
            filename = secure_filename(file.filename)
            filepath = os.path.join(app.config['UPLOAD_FOLDER'], filename)
            file.save(filepath)
            files_to_cleanup.append(filepath)
            files_to_process.append(filepath)
            print(f"処理対象に追加: {filename}")
        
        if not files_to_process:
            cleanup_uploads(files_to_cleanup)
            return jsonify({'error': '処理するPPTXファイルが見つかりません'}), 400
        
        # 複数ファイル処理の場合はZIPで返す
        if len(files_to_process) > 1:
            import zipfile
            zip_buffer = BytesIO()
            with zipfile.ZipFile(zip_buffer, 'w', zipfile.ZIP_DEFLATED) as zip_file:
                for file_path in files_to_process:
                    try:
                        prs = Presentation(file_path)
                        is_delete = (action == 'delete')
                        process_presentation(
                            prs, 
                            keywords, 
                            new_keyword if not is_delete else None,
                            is_delete=is_delete
                        )
                        
                        output = BytesIO()
                        prs.save(output)
                        output.seek(0)
                        
                        result_filename = f"modified_{os.path.basename(file_path)}"
                        zip_file.writestr(result_filename, output.getvalue())
                    except Exception as e:
                        print(f"ファイル処理エラー {file_path}: {str(e)}")
                        continue
            
            zip_buffer.seek(0)
            response = send_file(
                zip_buffer,
                mimetype='application/zip',
                as_attachment=True,
                download_name='modified_presentations.zip'
            )
            cleanup_uploads(files_to_cleanup)
            return response
        else:
            # 単一ファイル処理
            file_path = files_to_process[0]
            prs = Presentation(file_path)
            
            is_delete = (action == 'delete')
            modified_count = process_presentation(
                prs, 
                keywords, 
                new_keyword if not is_delete else None,
                is_delete=is_delete
            )
            
            output = BytesIO()
            prs.save(output)
            output.seek(0)
            
            result_filename = f"modified_{os.path.basename(file_path)}"
            
            response = send_file(
                output,
                mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation',
                as_attachment=True,
                download_name=result_filename
            )
            cleanup_uploads(files_to_cleanup)
            return response
    
    except Exception as e:
        cleanup_uploads(files_to_cleanup)
        return jsonify({'error': f'エラーが発生しました: {str(e)}'}), 500


@app.route('/api/preview', methods=['POST'])
def preview_results():
    """置換前後のプレビューAPI"""
    files_to_cleanup = []
    try:
        files = request.files.getlist('file')  # 複数ファイルに対応
        if not files or (len(files) == 1 and files[0].filename == ''):
            return jsonify({'error': 'ファイルがアップロードされていません'}), 400
        
        keywords_json = request.form.get('keywords', '[]')
        new_keyword = request.form.get('new_keyword', '').strip()
        action = request.form.get('action', 'replace')
        recursive = request.form.get('recursive', 'false').lower() == 'true'
        
        try:
            keywords = json.loads(keywords_json)
        except json.JSONDecodeError:
            keywords = [keywords_json] if keywords_json else []
        
        if not keywords or len(keywords) == 0:
            return jsonify({'error': 'キーワードを入力してください'}), 400
        
        if action == 'replace' and not new_keyword:
            return jsonify({'error': '置換先のキーワードを入力してください'}), 400
        
        # 複数ファイルを処理
        files_to_process = []
        for file in files:
            if file.filename == '':
                continue
            # 隠しファイルを除外
            if file.filename.startswith('.'):
                print(f"スキップ（隠しファイル）: {file.filename}")
                continue
            if not allowed_file(file.filename):
                print(f"スキップ（拡張子不可）: {file.filename}")
                continue
            
            filename = secure_filename(file.filename)
            filepath = os.path.join(app.config['UPLOAD_FOLDER'], filename)
            file.save(filepath)
            files_to_cleanup.append(filepath)
            files_to_process.append(filepath)
            print(f"処理対象に追加: {filename}")
        
        if not files_to_process:
            cleanup_uploads(files_to_cleanup)
            return jsonify({'error': '処理するPPTXファイルが見つかりません'}), 400
        
        # 全ファイルの統計を集約
        total_before_count = 0
        total_before_slides = 0
        total_after_count = 0
        total_after_slides = 0
        total_modified = 0
        
        for file_path in files_to_process:
            try:
                prs = Presentation(file_path)
                
                # 処理前の検出
                before_results = find_keywords_in_presentation(prs, keywords)
                before_count = sum(r['count'] for r in before_results)
                
                # 処理を実行（プレビューのみ）
                is_delete = (action == 'delete')
                modified_count = process_presentation(
                    prs, 
                    keywords, 
                    new_keyword if not is_delete else None,
                    is_delete=is_delete
                )
                
                # 処理後の検出
                after_results = find_keywords_in_presentation(prs, keywords)
                after_count = sum(r['count'] for r in after_results)
                
                total_before_count += before_count
                total_before_slides += len(before_results)
                total_after_count += after_count
                total_after_slides += len(after_results)
                total_modified += modified_count
            except Exception as e:
                print(f"ファイル処理エラー {file_path}: {str(e)}")
                continue
        
        response = jsonify({
            'success': True,
            'before': {
                'count': total_before_count,
                'slides': total_before_slides
            },
            'after': {
                'count': total_after_count,
                'slides': total_after_slides
            },
            'modified_shapes': total_modified,
            'files_processed': len(files_to_process),
            'action': action
        })
        
        cleanup_uploads(files_to_cleanup)
        return response
    
    except Exception as e:
        cleanup_uploads(files_to_cleanup)
        return jsonify({'error': f'エラーが発生しました: {str(e)}'}), 500


if __name__ == '__main__':
    app.run(debug=True, host='0.0.0.0', port=5000)
