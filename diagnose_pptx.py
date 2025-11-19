"""
PowerPoint ファイルの置換テスト
"""
from pptx import Presentation
import re
import os

def diagnose_pptx(filepath):
    """PowerPoint ファイルの内容を診断"""
    print(f"\n診断対象: {filepath}")
    print("=" * 60)
    
    if not os.path.exists(filepath):
        print(f"⚠️ ファイルが見つかりません: {filepath}")
        return
    
    try:
        prs = Presentation(filepath)
        print(f"スライド数: {len(prs.slides)}")
        
        total_keyword_count = 0
        
        for slide_num, slide in enumerate(prs.slides, 1):
            print(f"\n--- スライド {slide_num} ---")
            slide_has_content = False
            
            for shape_num, shape in enumerate(slide.shapes):
                if hasattr(shape, "text") and shape.text.strip():
                    slide_has_content = True
                    print(f"  [シェイプ {shape_num}]")
                    print(f"    テキスト: {shape.text}")
                    print(f"    テキストフレーム有無: {hasattr(shape, 'text_frame')}")
                    
                    if hasattr(shape, "text_frame"):
                        print(f"    段落数: {len(shape.text_frame.paragraphs)}")
                        for para_num, para in enumerate(shape.text_frame.paragraphs):
                            print(f"      [段落 {para_num}]")
                            print(f"        テキスト: {para.text}")
                            print(f"        ラン数: {len(para.runs)}")
                            for run_num, run in enumerate(para.runs):
                                print(f"          [ラン {run_num}] '{run.text}'")
            
            if not slide_has_content:
                print("  (テキストなし)")
        
        print("\n" + "=" * 60)
        
    except Exception as e:
        print(f"❌ エラー: {str(e)}")


def test_replace_in_pptx(input_filepath, keyword, new_keyword, output_filepath):
    """PowerPoint ファイル内のキーワードを置換してテスト"""
    print(f"\n置換テスト")
    print("=" * 60)
    print(f"入力ファイル: {input_filepath}")
    print(f"キーワード: '{keyword}'")
    print(f"置換先: '{new_keyword}'")
    print(f"出力ファイル: {output_filepath}")
    
    if not os.path.exists(input_filepath):
        print(f"❌ ファイルが見つかりません")
        return
    
    try:
        prs = Presentation(input_filepath)
        modified_count = 0
        total_replacements = 0
        
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text_frame"):
                    original_text = shape.text
                    
                    if keyword.lower() in original_text.lower():
                        print(f"\n  置換対象を発見:")
                        print(f"    前: {original_text}")
                        
                        # 置換実行
                        for paragraph in shape.text_frame.paragraphs:
                            full_text = ''.join(run.text for run in paragraph.runs)
                            if keyword.lower() in full_text.lower():
                                pattern = re.compile(re.escape(keyword), re.IGNORECASE)
                                new_full_text = pattern.sub(new_keyword, full_text)
                                count = pattern.subn(new_keyword, full_text)[1]
                                total_replacements += count
                                
                                # すべてのrunをクリア
                                for run in paragraph.runs:
                                    run.text = ''
                                
                                # 新しいテキストを設定
                                if paragraph.runs:
                                    paragraph.runs[0].text = new_full_text
                                else:
                                    paragraph.text = new_full_text
                        
                        new_text = shape.text
                        print(f"    後: {new_text}")
                        
                        if new_text != original_text:
                            modified_count += 1
        
        # ファイルを保存
        prs.save(output_filepath)
        
        print(f"\n✓ 完了")
        print(f"  修正されたシェイプ数: {modified_count}")
        print(f"  置換回数: {total_replacements}")
        print(f"  出力ファイル: {output_filepath}")
        
    except Exception as e:
        print(f"❌ エラー: {str(e)}")
        import traceback
        traceback.print_exc()


if __name__ == '__main__':
    # Test2.pptx のパスを設定
    test_file = 'd:\\genAI\\Tool\\uploads\\Test2.pptx'
    output_file = 'd:\\genAI\\Tool\\uploads\\Test2_replaced.pptx'
    
    # 診断
    diagnose_pptx(test_file)
    
    # 置換テスト
    test_replace_in_pptx(test_file, 'hitachi astemo', 'astemo', output_file)
    
    # 置換後のファイルを診断
    diagnose_pptx(output_file)
